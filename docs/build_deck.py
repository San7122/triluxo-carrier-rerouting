"""Generate docs/deck.pptx -- the Part 3 product-strategy presentation.

14 executive slides, built directly with python-pptx. Design intent: sparse
slides (a headline + one visual or <=4 bullets), depth carried in speaker notes.
EVERY number comes from the real evaluation run (eval/results.json); nothing is
invented. Regenerate:  python docs/build_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]

# ---- palette -------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x25, 0x45)
INK = RGBColor(0x1A, 0x1A, 0x2E)
TEAL = RGBColor(0x1B, 0x99, 0x8B)
AMBER = RGBColor(0xE8, 0xA3, 0x3D)
RED = RGBColor(0xD6, 0x45, 0x50)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREY = RGBColor(0x5B, 0x64, 0x70)
LIGHT = RGBColor(0xF1, 0xF4, 0xF7)
CLOUD = RGBColor(0xE9, 0xEF, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_num = [0]


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp, x, y, w, h)
    _fill(s, color)
    return s


def textbox(slide, x, y, w, h, lines, size=18, color=INK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
            space_after=8, line_spacing=1.05, gap=None):
    if gap is not None:
        space_after = gap
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        text, opts = (ln if isinstance(ln, tuple) else (ln, {}))
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.name = font
        f.color.rgb = opts.get("color", color)
    return tb


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(title, kicker=None):
    """Standard slide: left accent spine + kicker + title + underline."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, Inches(0.22), SH, TEAL)          # left accent spine
    if kicker:
        textbox(s, Inches(0.7), Inches(0.40), Inches(11.5), Inches(0.4),
                [(kicker.upper(), {"size": 13, "bold": True, "color": TEAL})])
    textbox(s, Inches(0.7), Inches(0.70), Inches(12.1), Inches(1.0),
            [(title, {"size": 27, "bold": True, "color": NAVY})])
    rect(s, Inches(0.72), Inches(1.62), Inches(2.1), Pt(3), AMBER)
    # slide number
    _num[0] += 1
    textbox(s, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
            [(str(_num[0]), {"size": 11, "color": GREY})], align=PP_ALIGN.RIGHT)
    return s


def bullets(slide, items, x=Inches(0.75), y=Inches(1.95), w=Inches(11.9),
            h=Inches(4.9), size=17, gap=13):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        text, opts = (it if isinstance(it, tuple) else (it, {}))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = ("•  " if opts.get("bullet", True) else "") + text
        run.font.size = Pt(opts.get("size", size))
        run.font.color.rgb = opts.get("color", INK)
        run.font.bold = opts.get("bold", False)
        run.font.name = "Calibri"
    return tb


def table(slide, x, y, col_w, rows_data, font=11, header_fill=NAVY,
          row_h=0.42, highlight=None):
    """Native table with explicit fills. col_w in EMU (use Inches()). `highlight`
    = dict {row_index: RGBColor} to tint a body row."""
    nrows, ncols = len(rows_data), len(rows_data[0])
    gf = slide.shapes.add_table(nrows, ncols, x, y, Emu(sum(col_w)),
                                Inches(row_h * nrows))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = Emu(w)
    for ri in range(nrows):
        tbl.rows[ri].height = Inches(row_h)
        for cij in range(ncols):
            cell = tbl.cell(ri, cij)
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            val = rows_data[ri][cij]
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if cij == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            f = run.font
            f.name = "Calibri"
            f.size = Pt(font)
            if ri == 0:
                f.bold = True
                f.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                f.color.rgb = INK
                cell.fill.solid()
                if highlight and ri in highlight:
                    cell.fill.fore_color.rgb = highlight[ri]
                else:
                    cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
    return gf


def flow_box(slide, x, y, w, h, title, sub, color, tcolor=WHITE):
    box = rect(slide, x, y, w, h, color, rounded=True)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = tcolor
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = tcolor
    return box


def arrow(slide, x, y, w, color=GREY, vertical=False):
    if vertical:
        rect(slide, x, y, Pt(3), w, color)
    else:
        rect(slide, x, y, w, Pt(3), color)


# =====================================================================
# 1 — Title
# =====================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(6.85), SW, Inches(0.65), TEAL)
textbox(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(1.5),
        [("Autonomous Carrier Rerouting", {"size": 44, "bold": True, "color": WHITE})])
textbox(s, Inches(0.92), Inches(3.05), Inches(11.6), Inches(1.0),
        [("A self-healing LangGraph agent — and a trajectory-based benchmark of "
          "when to trust it to act", {"size": 21, "color": RGBColor(0xC9, 0xE4, 0xE0)})])
textbox(s, Inches(0.95), Inches(4.55), Inches(11.4), Inches(0.5),
        [("Candidate: Sanjana Thakur   ·   AI Researcher Take-Home   ·   Triluxo Pvt. Ltd.",
          {"size": 15, "color": RGBColor(0x9F, 0xB3, 0xC8)})])
textbox(s, Inches(0.95), Inches(5.35), Inches(11.4), Inches(0.9),
        [("Stack: Python · LangGraph · OpenAI GPT-4o · Llama 3.3 70B & 3.1 8B (Groq) · "
          "pytest · Docker", {"size": 14, "color": RGBColor(0x8F, 0xA6, 0xBD)})])
notes(s, "In one line: I built a working autonomous rerouting agent, then evaluated a frontier "
         "closed model (GPT-4o) against two open models by scoring their decision-making step-by-step "
         "— not just outcomes — and I have a defensible, data-backed recommendation on which to trust "
         "in production. 15 minutes.")

# =====================================================================
# 2 — Business Problem
# =====================================================================
s = content_slide("A delayed shipment is a race against the clock — and humans are the bottleneck",
                  kicker="The business problem")
bullets(s, [
    ("A customs hold, port congestion, or carrier failure adds 14–40h to an ETA — "
     "on high-value, time-critical freight.", {}),
    ("Today a human expediter must notice, pull alternatives, weigh cost vs. speed "
     "vs. reliability, and rebook — often hours later, often off-shift.", {}),
    ("Every hour compounds: missed connections, penalty clauses, spoiled goods, SLA breaches.", {}),
    ("The decision is repetitive and rule-based — a strong fit for an autonomous agent, "
     "IF it can be trusted to take a real, irreversible action.", {"bold": True, "color": NAVY}),
], y=Inches(2.05))
notes(s, "Frame the pain around cost-of-delay and the human bottleneck. The punchline is the last "
         "bullet: this is automatable, but automation that books real freight and spends money must be "
         "trustworthy. That 'if it can be trusted' sets up the entire research contribution.")

# =====================================================================
# 3 — Research Objective
# =====================================================================
s = content_slide("The research question: which model can you trust to ACT?", kicker="Research objective")
bullets(s, [
    ("Assignment goal: compare a frontier closed-source model vs a leading open-source "
     "model on an agentic workflow, evaluating intermediate reasoning, tool-calling, and "
     "error recovery — not just the final output.", {}),
    ("Why trajectory evaluation: for an agent that takes irreversible actions, a right "
     "outcome reached by unsound reasoning — or sound reasoning followed by a wrong action "
     "— is a hidden liability. You must score the whole path.", {"bold": True, "color": NAVY}),
    ("RQ1: Can an open model match a frontier closed model on a deterministic logistics task?", {}),
    ("RQ2: Where do agents fail — and are those failures visible without inspecting the trajectory?", {}),
    ("RQ3: What must wrap the model for autonomous execution to be safe?", {}),
], y=Inches(1.95), gap=12)
notes(s, "State the three research questions explicitly — interviewers reward a crisp framing. The "
         "central methodological claim is RQ-independent: outcome-only testing is dangerous for "
         "action-taking agents. Everything downstream measures the trajectory to answer these.")

# =====================================================================
# 4 — System Architecture
# =====================================================================
s = content_slide("System architecture: an explicit LangGraph state machine", kicker="Architecture")
y0 = Inches(2.15)
flow_box(s, Inches(0.75), y0, Inches(2.7), Inches(1.15), "1 · Ingestion", "parse alert, rate severity", NAVY)
flow_box(s, Inches(3.85), y0, Inches(2.7), Inches(1.15), "2 · Evaluation", "call carrier API, weigh trade-offs", NAVY)
flow_box(s, Inches(6.95), y0, Inches(2.7), Inches(1.15), "3 · Decision", "pick per policy, book reroute", TEAL)
flow_box(s, Inches(10.05), y0, Inches(2.5), Inches(1.15), "✔ Rerouted", "booking confirmed", GREEN)
for ax in (Inches(3.5), Inches(6.6), Inches(9.7)):
    arrow(s, ax, Inches(2.63), Inches(0.35))
# fallback / escalation
fb = rect(s, Inches(3.85), Inches(4.05), Inches(8.65), Inches(1.0), AMBER, rounded=True)
tf = fb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = ("State transitions: tool fails OR no compliant option  →  retry once (budget = 1)  "
          "→  still failing  →  escalate to HUMAN REVIEW")
r.font.size = Pt(12.5)
r.font.bold = True
r.font.color.rgb = INK
arrow(s, Inches(5.0), Inches(3.3), Inches(0.75), GREY, vertical=True)
arrow(s, Inches(8.2), Inches(3.3), Inches(0.75), GREY, vertical=True)
textbox(s, Inches(0.75), Inches(5.45), Inches(12), Inches(1.3), [
    ("Why LangGraph: an explicit graph — not an open-ended ReAct loop — makes the recovery "
     "policy auditable and makes every node a scoreable trajectory step. Guardrail in node 3 "
     "rejects un-offered / non-compliant carriers BEFORE booking.", {"size": 14.5, "color": GREY}),
], anchor=MSO_ANCHOR.TOP)
notes(s, "Walk left to right: ingest -> evaluate (calls the simulated carrier API) -> decide & execute. "
         "The amber band is the safety net and the state-transition logic: any tool failure or "
         "no-viable-option situation retries once, then hands off to a human rather than crashing or "
         "guessing. Mermaid source of this graph is in the speaker-notes appendix and docs/research_report.md. "
         "The deterministic guardrail (nodes.py:184-196) is the single most important production element.")

# =====================================================================
# 5 — Agent Workflow (stages, honestly mapped)
# =====================================================================
s = content_slide("The rerouting pipeline — reference stages mapped to what's implemented",
                  kicker="Agent workflow")
stages = [
    ("Telemetry alert", "disruption JSON ingested", CLOUD, INK),
    ("Planner / Ingestion", "LLM: severity + reroute?", NAVY, WHITE),
    ("Risk + Route optimisation", "LLM: cost/ETA/reliability trade-off", NAVY, WHITE),
    ("Policy validation", "deterministic guardrail (enforced)", TEAL, WHITE),
    ("Execution", "LLM commits execute_reroute", NAVY, WHITE),
    ("Monitoring / Feedback", "trajectory logs + escalation check", GREY, WHITE),
]
x = Inches(0.7)
for i, (t, sub, col, tc) in enumerate(stages):
    flow_box(s, x, Inches(2.35), Inches(1.92), Inches(1.5), t, sub, col, tc)
    if i < len(stages) - 1:
        arrow(s, x + Inches(1.92), Inches(3.03), Inches(0.14))
    x += Inches(2.06)
textbox(s, Inches(0.72), Inches(4.15), Inches(12), Inches(2.4), [
    ("Implemented as LLM agents: Ingestion, Evaluation (risk+route), Execution — the parts worth "
     "evaluating.", {"size": 14, "color": INK}),
    ("Deterministic (not the LLM): Policy validation + the guardrail, and the retry-then-escalate "
     "control flow — safety must not depend on model whim.", {"size": 14, "color": TEAL, "bold": True}),
    ("Monitoring/Feedback today = structured trajectory logs + a post-hoc policy sanity-check on "
     "escalation; persistent memory & a live feedback loop are on the roadmap (slide 12).", {"size": 14, "color": GREY}),
], gap=11)
notes(s, "Be honest and precise here — a strong panel checks this. The classic control-tower stages "
         "(planner, risk, route optimiser, policy validation, execution, monitoring, feedback) are the "
         "REFERENCE pipeline. I implemented three of them as LLM agents (the parts where model judgment "
         "matters and is worth grading), kept policy validation and recovery DETERMINISTIC (safety), and "
         "scoped monitoring/feedback to trajectory logging for this 3-4h build. I deliberately did NOT "
         "spin up seven thin agents — that would be complexity without benefit.")

# =====================================================================
# 6 — Trajectory-Based Evaluation
# =====================================================================
s = content_slide("We grade the working, not just the final answer", kicker="Trajectory-based evaluation")
bullets(s, [
    ("Analogy: two students both write '42'. One reasoned it out; one guessed. Same answer, "
     "very different trust — you must see the working.", {"bold": True, "color": NAVY}),
    ("We log every step: the agent's reasoning, which tools it called with what arguments, and "
     "what it finally did.", {}),
    ("Real example from the logs — Llama 3.1 8B on 'no viable option': its analysis correctly "
     "flagged every carrier as non-compliant, then it BOOKED one anyway ($3,200 vs $3,000 cap; "
     "0.85 vs 0.90 reliability). Outcome-only scoring reads this as 'rerouted = success'.", {"color": RED}),
    ("Only the trajectory reveals the safety violation. That is the core contribution.", {"bold": True, "color": NAVY}),
], y=Inches(1.95), gap=13)
notes(s, "This is the heart of the technical contribution, explained for a mixed audience. The student "
         "analogy lands the point for non-technical stakeholders. The 8B example is the proof: a model can "
         "produce impeccable analysis and then take an unsafe action. If you only check the final state "
         "('a carrier was booked'), you miss it entirely. Three of our four scoring dimensions are computed "
         "deterministically against a policy oracle, so this is objective, not opinion.")

# =====================================================================
# 7 — Experimental Methodology
# =====================================================================
s = content_slide("Experimental methodology: identical harness, three models", kicker="Methodology")
table(s, Inches(0.75), Inches(2.0),
      [Inches(3.4), Inches(2.7), Inches(2.6), Inches(2.6)],
      [["Model", "Type", "Served via", "Role"],
       ["GPT-4o", "Closed frontier", "OpenAI API", "Closed baseline"],
       ["Llama 3.3 70B", "Open-source", "Groq", "Primary open candidate"],
       ["Llama 3.1 8B", "Open-source", "Groq", "Budget / speed option"]],
      font=12, row_h=0.5)
textbox(s, Inches(0.75), Inches(4.35), Inches(12), Inches(2.4), [
    ("4 scored dimensions (1–5): tool-calling accuracy · decision correctness · error recovery · "
     "reasoning quality (proxy, excluded from the headline number).", {"size": 14.5, "color": INK}),
    ("Ground truth = a deterministic policy ORACLE, so 3 of 4 dimensions are objective & reproducible.", {"size": 14.5, "color": TEAL, "bold": True}),
    ("5 scenarios stress distinct paths: normal · tight-margin trade-off · no-viable-option (must "
     "escalate) · transient tool failure (retry) · hard tool failure (retry-then-escalate). Temp 0.", {"size": 14.5, "color": INK}),
], gap=12)
notes(s, "The design that makes the comparison fair: identical graph, prompts, tools, and policy across "
         "all three models — only the client swaps. The oracle (policy_optimal) computes the correct answer "
         "independently, turning 'decision correctness' into a checkable number. The five scenarios are hand-"
         "built to exercise the exact behaviours that matter for autonomous action: optimal selection, "
         "discipline under ambiguity, knowing when to escalate, and recovering from tool failures.")

# =====================================================================
# 8 — Results
# =====================================================================
s = content_slide("Results: the open 70B matched the frontier model — and was safer than the 8B",
                  kicker="Results")
cd = CategoryChartData()
cd.categories = ["Tool use", "Decision\ncorrectness", "Error\nrecovery", "Overall"]
cd.add_series("Llama 3.3 70B (open)", (5.0, 5.0, 5.0, 5.0))
cd.add_series("GPT-4o (closed)", (5.0, 4.4, 4.33, 4.6))
cd.add_series("Llama 3.1 8B (open)", (4.4, 3.2, 3.67, 3.87))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.72), Inches(2.0), Inches(7.2), Inches(3.4), cd)
ch = gf.chart
ch.has_title = False
ch.value_axis.minimum_scale = 0
ch.value_axis.maximum_scale = 5
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.plots[0].series[0].format.fill.solid(); ch.plots[0].series[0].format.fill.fore_color.rgb = TEAL
ch.plots[0].series[1].format.fill.solid(); ch.plots[0].series[1].format.fill.fore_color.rgb = NAVY
ch.plots[0].series[2].format.fill.solid(); ch.plots[0].series[2].format.fill.fore_color.rgb = AMBER
# per-scenario table
table(s, Inches(0.72), Inches(5.55),
      [Inches(3.0), Inches(1.35), Inches(1.5), Inches(1.35)],
      [["Scenario", "GPT-4o", "70B", "8B"],
       ["tight_margin", "4.0", "5.0", "4.0"],
       ["no_viable_option", "5.0", "5.0", "1.33"]],
      font=10.5, row_h=0.34, highlight={2: RGBColor(0xFB, 0xE3, 0xE3)})
textbox(s, Inches(8.15), Inches(2.05), Inches(4.6), Inches(4.9), [
    ("Overall (obj. dims) · safety violations:", {"size": 13, "bold": True, "color": NAVY}),
    ("Llama 3.3 70B — 5.0 · 0 violations", {"size": 14, "color": GREEN, "bold": True}),
    ("GPT-4o — 4.6 · 0 violations", {"size": 14, "color": INK}),
    ("Llama 3.1 8B — 3.87 · 1 violation", {"size": 14, "color": RED}),
    ("", {"size": 6}),
    ("GPT-4o & the 8B made the SAME tight-margin slip: booked BAL-12, not the "
     "policy-optimal FAS-11.", {"size": 13, "color": INK}),
    ("Only the 8B forced a non-compliant booking (no_viable_option).", {"size": 13, "color": RED}),
    ("The open 70B: highest score, 0 violations, ~1/10th GPT-4o's per-token cost.", {"size": 13, "color": TEAL, "bold": True}),
], gap=8)
notes(s, "Read it honestly. GPT-4o is safe (0 violations) but NOT flawless — it lost decision points on "
         "the ambiguous tight-margin case, the exact same slip the 8B made (booking the safer-looking "
         "slower carrier instead of the policy-optimal fastest one). The open 70B was the single best "
         "performer with zero violations. The 8B is the only model that forced a policy-violating booking. "
         "Caveat I'll state up front: n=1 per scenario, so these are directional, not powered — the harness "
         "supports --trials for variance. Numbers are from eval/results.json.")

# =====================================================================
# 9 — Technical Trade-offs
# =====================================================================
s = content_slide("Technical trade-offs: open-source vs closed-source", kicker="Trade-offs")
table(s, Inches(0.75), Inches(2.05),
      [Inches(2.7), Inches(3.05), Inches(3.05), Inches(2.85)],
      [["Axis", "GPT-4o (closed)", "Llama 3.3 70B (open)", "Llama 3.1 8B (open)"],
       ["Overall score", "4.6", "5.0  (best)", "3.87"],
       ["Safety violations", "0", "0", "1"],
       ["Decision correctness", "4.4", "5.0", "3.2"],
       ["Cost / M tokens", "~$2.5 in / $10 out", "~$0.6–0.9", "~$0.05"],
       ["Tokens (5 runs)", "10,519 (leanest)", "17,206", "16,563"],
       ["Deployment", "Vendor API", "Self-host or hosted", "Self-host, cheap"]],
      font=11.5, row_h=0.52, highlight={2: RGBColor(0xE7, 0xF3, 0xEA)})
textbox(s, Inches(0.75), Inches(6.35), Inches(12), Inches(0.9), [
    ("Latency is cross-provider (OpenAI vs Groq's accelerated, rate-limit-paced inference) — "
     "directional only, not an SLA. Accuracy & safety are the decisive axes here.", {"size": 13, "color": GREY}),
], anchor=MSO_ANCHOR.TOP)
notes(s, "The honest trade-off: on THIS narrow, well-specified task, the open 70B wins on accuracy and "
         "matches on safety, at roughly a tenth of the cost, with the option to self-host (data residency, "
         "no per-call vendor dependency). GPT-4o's advantages — broader capability, less prompt fragility on "
         "open-ended tasks — don't get exercised by a deterministic policy workflow. I am NOT claiming open "
         "beats closed in general; I'm claiming it's the right tool for this job.")

# =====================================================================
# 10 — Production Readiness
# =====================================================================
s = content_slide("Production readiness: built like a system, not a notebook", kicker="Production readiness")
items = [
    ("Guardrails", "deterministic block on un-offered / non-compliant bookings", TEAL),
    ("Error recovery", "retry-once-then-escalate; never fakes success", TEAL),
    ("Testing", "12 unit tests incl. one reproducing the 8B failure scores", NAVY),
    ("Reproducibility", "rescore from committed logs, no API key (eval.score)", NAVY),
    ("Docker", "one-command build + key-free smoke test", GREY),
    ("Logging & monitoring", "structured control-flow logs; process-mining view", GREY),
]
x, y = Inches(0.75), Inches(2.15)
for i, (t, sub, col) in enumerate(items):
    col_i = i % 2
    row_i = i // 2
    bx = Inches(0.75) + col_i * Inches(6.15)
    by = Inches(2.15) + row_i * Inches(1.45)
    card = rect(s, bx, by, Inches(5.85), Inches(1.25), LIGHT, rounded=True)
    rect(s, bx, by, Inches(0.12), Inches(1.25), col)
    textbox(s, bx + Inches(0.3), by + Inches(0.12), Inches(5.4), Inches(1.0),
            [(t, {"size": 15, "bold": True, "color": NAVY}),
             (sub, {"size": 12.5, "color": GREY})], gap=3)
textbox(s, Inches(0.75), Inches(6.65), Inches(12), Inches(0.6), [
    ("Scalability path: thread-safe rate limiter; swap simulated tools for TMS/booking APIs "
     "behind the same interfaces; LangGraph checkpointing for crash-safe in-flight reroutes.", {"size": 12.5, "color": GREY}),
], anchor=MSO_ANCHOR.TOP)
notes(s, "This is the engineering-maturity slide. The message: safety is DETERMINISTIC and does not depend "
         "on the model; the results are reproducible WITHOUT an API key (I rescored from committed "
         "trajectories); and the whole thing runs in Docker with one command. The test that reproduces the "
         "8B's exact failure scores is the credibility anchor — it proves the benchmark itself is correct.")

# =====================================================================
# 11 — Business Recommendation
# =====================================================================
s = content_slide("Can an open model replace GPT-4o here? For this task — yes", kicker="Business recommendation")
c1 = rect(s, Inches(0.75), Inches(2.1), Inches(5.75), Inches(3.5), RGBColor(0xEA, 0xF6, 0xF4), rounded=True)
c2 = rect(s, Inches(6.8), Inches(2.1), Inches(5.75), Inches(3.5), RGBColor(0xFB, 0xF0, 0xDE), rounded=True)
textbox(s, Inches(1.0), Inches(2.35), Inches(5.3), Inches(3.1), [
    ("✔ Deploy Llama 3.3 70B — recommended primary", {"size": 16, "bold": True, "color": RGBColor(0x1B, 0x77, 0x6B)}),
    ("Beat GPT-4o (5.0 vs 4.6), 0 safety violations, ~1/10th the cost.", {"size": 13.5, "color": INK}),
    ("Runs BEHIND the deterministic policy guardrail.", {"size": 13.5, "color": INK}),
    ("Self-hostable → data residency + no per-call vendor lock-in.", {"size": 13.5, "color": INK}),
], gap=9)
textbox(s, Inches(7.05), Inches(2.35), Inches(5.3), Inches(3.1), [
    ("✘ Do NOT let Llama 3.1 8B execute autonomously", {"size": 16, "bold": True, "color": RGBColor(0xB2, 0x38, 0x42)}),
    ("It forced a policy-violating booking — a safety defect.", {"size": 13.5, "color": INK}),
    ("Use only for read-only triage (ingest / summarise).", {"size": 13.5, "color": INK}),
    ("Keep GPT-4o as a swap-in reference / open-ended fallback.", {"size": 13.5, "color": INK}),
], gap=9)
textbox(s, Inches(0.75), Inches(5.85), Inches(12), Inches(1.1), [
    ("ROI: inference is ~$0.003 per reroute (70B) vs ~$15–30 of human expediter labour + hours of "
     "delay. The decision is about TRUST, not token price — and the guardrailed 70B earns it.", {"size": 14, "bold": True, "color": NAVY}),
], anchor=MSO_ANCHOR.TOP)
notes(s, "Give the direct answer the brief demands, no hedge. YES — the open 70B can replace GPT-4o for "
         "this use case, wrapped in the guardrails we already built, and it's cheaper and self-hostable. "
         "NO — the small 8B must not execute autonomously; it's cheap triage only. ROI framing for the "
         "non-technical stakeholders: inference cost is a rounding error next to human labour and cost-of-"
         "delay, so choose on safety, not price.")

# =====================================================================
# 12 — Future Improvements
# =====================================================================
s = content_slide("Future improvements", kicker="Roadmap")
bullets(s, [
    ("Broaden the closed baseline: run Claude & Gemini through the identical harness (adapters "
     "wired) to test cross-vendor generalisation.", {}),
    ("Statistical power: multi-trial (--trials) + 25+ scenarios + adversarial cases; LLM-as-judge "
     "to replace the keyword reasoning proxy.", {}),
    ("Multi-agent memory & human-in-the-loop: LangGraph checkpointing for crash-safe in-flight "
     "reroutes; an escalation console for the human reviewer.", {}),
    ("Observability at scale: process-mining dashboards over trajectory logs (path conformance, "
     "escalation & loop rates).", {}),
    ("Real integrations via MCP / TMS + booking-API tools, replacing the simulated carrier network "
     "behind the same interfaces.", {}),
], y=Inches(2.0), gap=13)
notes(s, "A credible, staged path. The single highest-value next step is broadening the closed baseline "
         "(Claude/Gemini) — the code already supports it. Then rigour (multi-trial, more scenarios, LLM-as-"
         "judge), then production scaffolding (memory, HITL console, MCP integrations, process-mining "
         "observability). Nothing here is hand-waving — each maps to a hook already in the codebase.")

# =====================================================================
# 13 — Lessons Learned
# =====================================================================
s = content_slide("Lessons learned", kicker="Reflection")
col_w = Inches(3.9)
titles = [("Research", TEAL), ("Engineering", NAVY), ("Product", AMBER)]
lessons = [
    ["Outcome-only eval hides safety defects — a model can reason correctly then act unsafely.",
     "The frontier model isn't automatically best on a narrow, well-specified task."],
    ["Make safety deterministic, not model-dependent — guardrails + explicit control flow.",
     "An independent oracle turns 'evaluation' from opinion into reproducible numbers."],
    ["Choose the model on trust, not token price — inference cost is a rounding error.",
     "'Where does the LLM add value?' — scope it to the ambiguous parts, not solved ones."],
]
for i, (t, col) in enumerate(titles):
    x = Inches(0.75) + i * Inches(4.15)
    rect(s, x, Inches(2.15), col_w, Inches(0.62), col, rounded=True)
    textbox(s, x, Inches(2.24), col_w, Inches(0.5), [(t, {"size": 16, "bold": True, "color": WHITE})],
            align=PP_ALIGN.CENTER)
    tb = s.shapes.add_textbox(x, Inches(3.0), col_w, Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, ln in enumerate(lessons[i]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = "•  " + ln
        r.font.size = Pt(13)
        r.font.color.rgb = INK
        r.font.name = "Calibri"
notes(s, "Three lenses. Research: the trajectory-vs-outcome insight, and the counter-intuitive result that "
         "a frontier model didn't dominate. Engineering: safety belongs in deterministic code, and an oracle "
         "makes evaluation objective. Product: cost is not the deciding variable — trust is — and the "
         "discipline of asking where the LLM actually earns its place.")

# =====================================================================
# 14 — Q&A Backup
# =====================================================================
s = content_slide("Q&A — anticipated questions", kicker="Appendix")
qa = [
    ("Isn't n=1 per scenario too thin?",
     "Yes — directional, not powered. The harness supports --trials; conclusions are hypotheses to widen."),
    ("Are the tools real?",
     "No — the carrier API is simulated & clearly labelled; interfaces match real TMS/booking calls for a localised swap."),
    ("Why did the open 70B beat GPT-4o?",
     "The task is a deterministic policy; GPT-4o's slip was one ambiguous trade-off, not a safety failure. Not a general claim."),
    ("Why an LLM if a rule can pick the carrier?",
     "For selection, the oracle can. The LLM earns its place on severity triage, ambiguity, and free-text — see slide 5/11."),
    ("How do I trust the benchmark itself?",
     "3 of 4 dims are deterministic vs an oracle; a unit test reproduces the 8B's exact failure scores; rescore needs no key."),
]
tb = s.shapes.add_textbox(Inches(0.75), Inches(1.95), Inches(11.9), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
for i, (q, a) in enumerate(qa):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run(); r.text = "Q:  " + q
    r.font.size = Pt(14.5); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.space_after = Pt(13)
    r2 = p2.add_run(); r2.text = "A:  " + a
    r2.font.size = Pt(13.5); r2.font.color.rgb = INK; r2.font.name = "Calibri"
notes(s, "Backup slide — only shown if asked. These are the five questions a sharp panel asks. The meta-"
         "message across all answers: I know the limits of my own work and can defend the methodology. That "
         "honesty is the point.")

out = ROOT / "docs" / "deck.pptx"
prs.save(str(out))
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
